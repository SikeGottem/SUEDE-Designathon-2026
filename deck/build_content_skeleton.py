#!/usr/bin/env python3
# Renders the evidence-led deck spec as an editable, brand-free review presentation.
"""Render the SUEDE deck spec as an editable, brand-free PPTX content skeleton."""

from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


ROOT = Path(__file__).resolve().parents[1]
DEFAULT_SPEC = ROOT / "deck" / "deck-spec.json"
DEFAULT_OUTPUT = ROOT / "deck" / "generated" / "SUEDE-2026-Content-Skeleton.pptx"

BLACK = RGBColor(0, 0, 0)
MUTED = RGBColor(102, 102, 102)
RULE = RGBColor(217, 217, 217)
WHITE = RGBColor(255, 255, 255)
FONT = "Arial"


def resolve_path(path: Path) -> Path:
    return path if path.is_absolute() else ROOT / path


def require_string(value: Any, label: str, errors: list[str]) -> str:
    if not isinstance(value, str):
        errors.append(f"{label} must be a string")
        return ""
    return value.strip()


def validate_spec(spec: Any) -> list[str]:
    """Check only the fields this renderer needs, without claiming final readiness."""
    errors: list[str] = []
    if not isinstance(spec, dict):
        return ["deck spec root must be an object"]
    deck = spec.get("deck")
    if not isinstance(deck, dict):
        errors.append("deck must be an object")
    elif deck.get("aspect_ratio") != "16:9":
        errors.append("deck.aspect_ratio must be 16:9")
    slides = spec.get("slides")
    if not isinstance(slides, list) or not slides:
        return errors + ["slides must be a non-empty array"]
    for number, slide in enumerate(slides, start=1):
        prefix = f"slides[{number - 1}]"
        if not isinstance(slide, dict):
            errors.append(f"{prefix} must be an object")
            continue
        for field in ("id", "section", "status", "audience_takeaway"):
            require_string(slide.get(field), f"{prefix}.{field}", errors)
        claim = slide.get("claim", "")
        if not isinstance(claim, str):
            errors.append(f"{prefix}.claim must be a string when provided")
        refs = slide.get("evidence_refs", [])
        if not isinstance(refs, list) or not all(isinstance(ref, str) for ref in refs):
            errors.append(f"{prefix}.evidence_refs must be an array of strings")
        visual = slide.get("visual")
        if not isinstance(visual, dict):
            errors.append(f"{prefix}.visual must be an object")
        elif not isinstance(visual.get("composition"), str):
            errors.append(f"{prefix}.visual.composition must be a string")
        notes = slide.get("speaker_notes", [])
        if not isinstance(notes, list) or not all(isinstance(note, str) for note in notes):
            errors.append(f"{prefix}.speaker_notes must be an array of strings")
        if not isinstance(slide.get("duration_seconds"), int):
            errors.append(f"{prefix}.duration_seconds must be an integer")
    return errors


def set_text_style(shape: Any, *, size: float, color: RGBColor, bold: bool = False,
                   align: PP_ALIGN = PP_ALIGN.LEFT) -> None:
    text_frame = shape.text_frame
    text_frame.margin_left = 0
    text_frame.margin_right = 0
    text_frame.margin_top = 0
    text_frame.margin_bottom = 0
    text_frame.vertical_anchor = MSO_ANCHOR.TOP
    for paragraph in text_frame.paragraphs:
        paragraph.alignment = align
        paragraph.space_after = Pt(0)
        for run in paragraph.runs:
            run.font.name = FONT
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color


def add_text(slide: Any, text: str, x: float, y: float, w: float, h: float, *,
             size: float, color: RGBColor = BLACK, bold: bool = False,
             name: str, align: PP_ALIGN = PP_ALIGN.LEFT) -> Any:
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    shape.name = name
    shape.text_frame.word_wrap = True
    shape.text_frame.text = text
    set_text_style(shape, size=size, color=color, bold=bold, align=align)
    return shape


def add_rule(slide: Any, x: float, y: float, w: float, *, name: str) -> None:
    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.012))
    rule.name = name
    rule.fill.solid()
    rule.fill.fore_color.rgb = RULE
    rule.line.fill.background()


def notes_text(slide_spec: dict[str, Any]) -> str:
    duration = slide_spec.get("duration_seconds", 0)
    prompts = [f"Timing: {duration} seconds."]
    prompts.extend(
        note.strip()
        for note in slide_spec.get("speaker_notes", [])
        if note.strip() and not note.strip().lower().startswith("timing:")
    )
    return "\n".join(prompts)


def write_notes(slide: Any, content: str) -> bool:
    """Use python-pptx notes support when provided by the installed version."""
    try:
        notes_slide = slide.notes_slide
        text_frame = getattr(notes_slide, "notes_text_frame", None)
        if text_frame is None:
            return False
        text_frame.text = content
        return True
    except (AttributeError, NotImplementedError):
        return False


def blocked_claim(slide_spec: dict[str, Any]) -> str:
    claim = slide_spec.get("claim", "").strip()
    if slide_spec.get("status") == "blocked" and claim:
        return f"BLOCKED — {claim}"
    if claim:
        return claim
    section = slide_spec.get("section", "this slide").lower()
    return f"BLOCKED — missing claim; resolve the {section} decision before final design."


def render_slide(prs: Presentation, slide_spec: dict[str, Any], index: int, total: int) -> bool:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = WHITE
    slide_id = slide_spec["id"]
    status = slide_spec["status"].upper()
    header = f"{index:02d}/{total:02d}   {slide_spec['section'].upper()}   /   {status}"
    add_text(slide, header, 0.75, 0.55, 11.85, 0.28, size=11, color=MUTED,
             bold=True, name=f"{slide_id}/index-section-status")
    add_rule(slide, 0.75, 0.97, 11.85, name=f"{slide_id}/header-rule")

    claim = blocked_claim(slide_spec)
    claim_color = BLACK if slide_spec["status"] != "blocked" else MUTED
    add_text(slide, claim, 0.75, 1.27, 11.85, 1.42, size=28, color=claim_color,
             bold=True, name=f"{slide_id}/claim-or-blocker")

    add_text(slide, "AUDIENCE TAKEAWAY", 0.75, 3.04, 2.0, 0.24, size=10, color=MUTED,
             bold=True, name=f"{slide_id}/takeaway-label")
    add_text(slide, slide_spec["audience_takeaway"], 0.75, 3.36, 7.4, 0.93, size=17,
             name=f"{slide_id}/audience-takeaway")

    visual = slide_spec.get("visual", {})
    add_text(slide, "COMPOSITION NOTE", 8.65, 3.04, 3.95, 0.24, size=10, color=MUTED,
             bold=True, name=f"{slide_id}/composition-label")
    add_text(slide, visual.get("composition", "No composition note supplied."), 8.65, 3.36, 3.95, 1.35,
             size=14, color=BLACK, name=f"{slide_id}/composition-note")

    add_rule(slide, 0.75, 5.14, 11.85, name=f"{slide_id}/footer-rule")
    refs = slide_spec.get("evidence_refs", [])
    ref_text = " • ".join(refs) if refs else "No evidence references supplied."
    add_text(slide, "EVIDENCE REFS", 0.75, 5.45, 1.7, 0.24, size=10, color=MUTED,
             bold=True, name=f"{slide_id}/evidence-label")
    add_text(slide, ref_text, 0.75, 5.75, 8.15, 0.72, size=12, color=MUTED,
             name=f"{slide_id}/evidence-refs")

    concise_notes = " ".join(
        note.strip()
        for note in slide_spec.get("speaker_notes", [])
        if note.strip() and not note.strip().lower().startswith("timing:")
    )
    timing = f"{slide_spec.get('duration_seconds', 0)} SEC"
    add_text(slide, timing, 10.65, 5.45, 1.95, 0.24, size=10, color=MUTED, bold=True,
             align=PP_ALIGN.RIGHT, name=f"{slide_id}/timing")
    add_text(slide, concise_notes or "No speaker prompt supplied.", 9.2, 5.75, 3.4, 0.72, size=11,
             color=MUTED, align=PP_ALIGN.RIGHT, name=f"{slide_id}/concise-notes")
    return write_notes(slide, notes_text(slide_spec))


def build(spec: dict[str, Any], output: Path) -> tuple[int, bool]:
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    notes_supported = True
    slides = spec["slides"]
    for index, slide_spec in enumerate(slides, start=1):
        notes_supported = render_slide(prs, slide_spec, index, len(slides)) and notes_supported
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)
    return len(slides), notes_supported


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--spec", type=Path, default=DEFAULT_SPEC,
                        help="Deck spec JSON (default: deck/deck-spec.json)")
    parser.add_argument("--output", type=Path, default=DEFAULT_OUTPUT,
                        help="PPTX output (default: deck/generated/SUEDE-2026-Content-Skeleton.pptx)")
    args = parser.parse_args()
    spec_path = resolve_path(args.spec)
    output = resolve_path(args.output)
    try:
        spec = json.loads(spec_path.read_text(encoding="utf-8"))
    except FileNotFoundError:
        print(f"ERROR: spec not found: {spec_path}", file=sys.stderr)
        return 2
    except json.JSONDecodeError as error:
        print(f"ERROR: invalid JSON in {spec_path}: {error}", file=sys.stderr)
        return 2
    except OSError as error:
        print(f"ERROR: could not read {spec_path}: {error}", file=sys.stderr)
        return 2
    errors = validate_spec(spec)
    if errors:
        for error in errors:
            print(f"ERROR: {error}", file=sys.stderr)
        return 1
    try:
        count, notes_supported = build(spec, output)
    except OSError as error:
        print(f"ERROR: could not write {output}: {error}", file=sys.stderr)
        return 2
    print(f"Built {count} editable slides: {output}")
    if notes_supported:
        print("Speaker notes written through python-pptx notes support.")
    else:
        print("Speaker-note limitation: installed python-pptx does not expose a writable notes text frame; concise notes remain on-slide.")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
