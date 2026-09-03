#!/usr/bin/env python3
"""Build Ethan's visual proposal deck for the SUEDE Designathon winning plan."""

from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
OUT = ROOT / "SUEDE-Designathon-2026-Winning-Plan.pptx"

W, H = Inches(13.333333), Inches(7.5)
BG = RGBColor(13, 10, 18)
PANEL = RGBColor(27, 22, 36)
PANEL_2 = RGBColor(38, 31, 50)
WHITE = RGBColor(247, 244, 250)
MUTED = RGBColor(183, 174, 196)
PURPLE = RGBColor(162, 105, 255)
LILAC = RGBColor(211, 180, 255)
LIME = RGBColor(201, 241, 116)
CORAL = RGBColor(255, 127, 111)
CYAN = RGBColor(105, 220, 230)


def rect(slide, x, y, w, h, color, radius=False, line=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = line if line else color
    if radius:
        shape.adjustments[0] = 0.12
    return shape


def text(slide, value, x, y, w, h, size=20, color=WHITE, bold=False,
         font="Avenir Next", align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP,
         margin=0, italic=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = value
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def rich(slide, runs, x, y, w, h, size=20, color=WHITE, font="Avenir Next",
         align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, spacing=1.0):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = spacing
    for item in runs:
        r = p.add_run()
        r.text = item[0]
        r.font.name = font
        r.font.size = Pt(item[1] if len(item) > 1 and item[1] else size)
        r.font.bold = bool(item[2]) if len(item) > 2 else False
        r.font.color.rgb = item[3] if len(item) > 3 and item[3] else color
    return box


def add_bg(slide):
    rect(slide, 0, 0, 13.333333, 7.5, BG)
    rect(slide, 0, 0, 0.12, 7.5, PURPLE)


def add_num(slide, n):
    text(slide, f"{n:02}", 12.25, 0.36, 0.55, 0.3, 10, MUTED, True, align=PP_ALIGN.RIGHT)


def title(slide, kicker, heading, n):
    text(slide, kicker.upper(), 0.62, 0.42, 6.5, 0.28, 10, PURPLE, True)
    text(slide, heading, 0.62, 0.78, 11.7, 0.68, 29, WHITE, True)
    add_num(slide, n)


def pill(slide, label, x, y, w, color=PURPLE, text_color=BG):
    rect(slide, x, y, w, 0.34, color, radius=True)
    text(slide, label.upper(), x, y + 0.02, w, 0.25, 9, text_color, True,
         align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)


def card(slide, x, y, w, h, eyebrow, heading, body, accent=PURPLE):
    rect(slide, x, y, w, h, PANEL, radius=True, line=PANEL_2)
    rect(slide, x, y, 0.08, h, accent, radius=True)
    text(slide, eyebrow.upper(), x + 0.28, y + 0.24, w - 0.55, 0.22, 9, accent, True)
    text(slide, heading, x + 0.28, y + 0.57, w - 0.55, 0.45, 18, WHITE, True)
    text(slide, body, x + 0.28, y + 1.12, w - 0.55, h - 1.3, 12.5, MUTED)


def bullet_list(slide, items, x, y, w, h, size=15, accent=PURPLE, gap=0.58):
    for i, item in enumerate(items):
        yy = y + i * gap
        rect(slide, x, yy + 0.09, 0.12, 0.12, accent, radius=True)
        text(slide, item, x + 0.28, yy, w - 0.28, gap, size, WHITE)


prs = Presentation()
prs.slide_width = W
prs.slide_height = H
blank = prs.slide_layouts[6]

# 1 — cover
s = prs.slides.add_slide(blank)
add_bg(s)
rect(s, 8.9, 0, 4.433333, 7.5, PURPLE)
text(s, "SUEDE DESIGNATHON 2026", 0.68, 0.66, 5.0, 0.3, 11, PURPLE, True)
rich(s, [("WIN\n", 42, True, WHITE), ("THE\n", 42, True, WHITE), ("WEEKEND.", 42, True, LIME)],
     0.68, 1.22, 7.4, 2.55, valign=MSO_ANCHOR.MIDDLE)
text(s, "A competition operating plan for", 0.72, 4.30, 5.8, 0.28, 15, MUTED)
text(s, "Ethan · Darius · CC · Akari · Chloe", 0.72, 4.72, 7.2, 0.36, 20, WHITE, True)
pill(s, "proposal · 1 sep 2026", 0.72, 5.58, 2.15, LILAC)
text(s, "3", 9.34, 0.63, 3.3, 1.32, 72, BG, True, align=PP_ALIGN.CENTER)
text(s, "days to prepare", 9.38, 1.83, 3.25, 0.35, 16, BG, True, align=PP_ALIGN.CENTER)
rect(s, 9.56, 2.60, 2.92, 0.035, BG)
text(s, "5", 9.34, 3.05, 3.3, 1.3, 72, BG, True, align=PP_ALIGN.CENTER)
text(s, "people · maximum team", 9.38, 4.20, 3.25, 0.65, 15, BG, True, align=PP_ALIGN.CENTER)
text(s, "One team. One loop. One story.", 9.35, 6.40, 3.25, 0.5, 14, BG, True, align=PP_ALIGN.CENTER)

# 2 — reality
s = prs.slides.add_slide(blank)
add_bg(s); title(s, "the reality", "This is a 50-hour decision race.", 2)
card(s, 0.65, 1.75, 3.72, 2.05, "Fri 4 Sep", "5–8pm", "Brief, rules, research, and concept lock. Leave with one direction.", PURPLE)
card(s, 4.80, 1.75, 3.72, 2.05, "Sat 5 Sep", "8am–8pm", "Prove the core loop, build it, test it, then freeze scope by 5pm.", CYAN)
card(s, 8.95, 1.75, 3.72, 2.05, "Sun 6 Sep", "8am–8pm", "Make the story undeniable. Submit early. Rehearse until calm.", LIME)
pill(s, "confirmed", 0.66, 4.32, 1.22, LIME)
text(s, "In person · Wilkinson Building G04 · USYD · university students · teams of 3–5", 2.08, 4.34, 10.0, 0.35, 14, WHITE, True)
pill(s, "unknown", 0.66, 5.14, 1.22, CORAL)
text(s, "2026 brief · rubric weights · pitch length · submission deadline · AI/IP rules · prizes", 2.08, 5.16, 10.2, 0.35, 14, WHITE, True)
text(s, "Rule: prepare the machine now. Do not pre-build the solution.", 0.66, 6.22, 11.6, 0.45, 22, LILAC, True)

# 3 — winner pattern
s = prs.slides.add_slide(blank)
add_bg(s); title(s, "what past winners reveal", "The winning shape is narrow, native, and retellable.", 3)
card(s, 0.65, 1.72, 3.75, 3.25, "2025 · 1st", "HOT TOPIC", "A redesigned comment experience under videos: one familiar behaviour, three prompts, a clear interaction, and no new destination app.", PURPLE)
card(s, 4.78, 1.72, 3.75, 3.25, "2025 · 2nd", "LLAMALERT", "A memorable llama flags misinformation, explains why, then turns learning into a tiny game. Instantly repeatable.", CORAL)
card(s, 8.91, 1.72, 3.75, 3.25, "2025 · 3rd", "TOUCAN", "A lightweight overlay at one moment of need. The team deliberately narrowed to learning, awareness, and convenience.", CYAN)
rich(s, [("INFERENCE  ", 12, True, PURPLE), ("The judges rewarded a precise intervention in an existing routine — backed by proof and presented as a story.", 18, True, WHITE)],
     0.68, 5.56, 11.7, 0.75)
text(s, "Not another all-in-one app.", 0.68, 6.48, 11.4, 0.35, 16, CORAL, True)

# 4 — equation
s = prs.slides.add_slide(blank)
add_bg(s); title(s, "our doctrine", "One user. One moment. One hero loop. One story.", 4)
labels = [
    ("USER", "Specific — never ‘everyone’", PURPLE),
    ("MOMENT", "A real trigger with pain", CYAN),
    ("LOOP", "Trigger → action → payoff", LIME),
    ("PROOF", "Observed behaviour, not vibes", CORAL),
    ("STORY", "Repeatable after one sentence", LILAC),
]
for i, (head, body, c) in enumerate(labels):
    x = 0.68 + i * 2.48
    rect(s, x, 1.84, 2.15, 3.42, PANEL, radius=True, line=PANEL_2)
    text(s, f"0{i+1}", x + 0.20, 2.05, 0.55, 0.3, 11, c, True)
    text(s, head, x + 0.20, 2.58, 1.75, 0.35, 17, WHITE, True)
    text(s, body, x + 0.20, 3.18, 1.75, 1.15, 13, MUTED)
    if i < 4:
        text(s, "+", x + 2.17, 3.17, 0.28, 0.5, 24, MUTED, True, align=PP_ALIGN.CENTER)
rect(s, 0.68, 5.74, 11.95, 0.78, PURPLE, radius=True)
text(s, "If the idea needs a paragraph before the demo, it is not ready.", 0.95, 5.91, 11.4, 0.40, 20, BG, True, align=PP_ALIGN.CENTER)

# 5 — scorecard
s = prs.slides.add_slide(blank)
add_bg(s); title(s, "concept selection", "No concept reaches high-fi below 20 / 25.", 5)
score = [
    ("Problem", "specific + evidenced", PURPLE),
    ("Effectiveness", "changes behaviour", CYAN),
    ("Innovation", "distinct mechanism", LIME),
    ("Craft", "excellent in time", CORAL),
    ("Story", "instant recall", LILAC),
]
for i, (head, sub, c) in enumerate(score):
    y = 1.72 + i * 0.88
    text(s, f"0{i+1}", 0.72, y + 0.06, 0.35, 0.28, 10, c, True)
    text(s, head, 1.24, y, 2.12, 0.35, 16, WHITE, True)
    text(s, sub, 3.48, y + 0.02, 2.15, 0.35, 13, MUTED)
    for j in range(5):
        rect(s, 6.38 + j * 1.05, y - 0.04, 0.66, 0.56, PANEL_2 if j < 4 else c, radius=True)
        text(s, str(j + 1), 6.38 + j * 1.05, y + 0.07, 0.66, 0.24, 12, BG if j == 4 else MUTED, True, align=PP_ALIGN.CENTER)
pill(s, "historic rubric", 0.72, 6.33, 1.72, LILAC)
text(s, "SUEDE 2024 scored visual communication, problem, solution, innovation, and presentation.", 2.66, 6.33, 9.65, 0.35, 13.5, WHITE)

# 6 — roles
s = prs.slides.add_slide(blank)
add_bg(s); title(s, "five owners · one team", "Clear accountability beats polite overlap.", 6)
roles = [
    ("ETHAN", "CAPTAIN / PRODUCT", "brief · scope · decisions · pitch", PURPLE),
    ("DARIUS", "RESEARCH / PROOF", "interviews · tests · evidence", CYAN),
    ("CC", "UX / INTERACTION", "flow · states · prototype logic", LIME),
    ("AKARI", "VISUAL / BRAND", "system · craft · motion", CORAL),
    ("CHLOE", "STORY / SUBMISSION", "deck · QA · backups · timing", LILAC),
]
for i, (name, role, job, c) in enumerate(roles):
    y = 1.62 + i * 0.98
    rect(s, 0.70, y, 11.95, 0.72, PANEL, radius=True, line=PANEL_2)
    rect(s, 0.70, y, 0.10, 0.72, c, radius=True)
    text(s, name, 1.02, y + 0.18, 1.35, 0.25, 13, c, True)
    text(s, role, 2.65, y + 0.17, 2.55, 0.26, 13, WHITE, True)
    text(s, job, 5.45, y + 0.17, 6.65, 0.28, 13, MUTED)
text(s, "Swap names after a 20-minute skills audit — keep these five ownership areas.", 0.72, 6.67, 11.6, 0.3, 13, MUTED, italic=True)

# 7 — pre-event
s = prs.slides.add_slide(blank)
add_bg(s); title(s, "1–3 september", "Prepare the operating system — legally.", 7)
card(s, 0.66, 1.66, 3.73, 4.82, "Tue · align", "Register + assign", "• confirm 5 registrations\n• finalise role owners\n• shared Figma + drive\n• empty research/deck templates\n• ask organisers 4 rule questions", PURPLE)
card(s, 4.80, 1.66, 3.73, 4.82, "Wed · rehearse", "Run the machine", "• 90-minute neutral mini-sprint\n• concept scorecard drill\n• five-minute interview\n• lo-fi loop\n• 60-second pitch\n• test file permissions", CYAN)
card(s, 8.94, 1.66, 3.73, 4.82, "Thu · de-risk", "Lower friction", "• pack + charge everything\n• fonts/tools available offline\n• travel + 4:15pm arrival\n• backup file system\n• stop heavy work by 9pm\n• sleep", LIME)

# 8 — Friday
s = prs.slides.add_slide(blank)
add_bg(s); title(s, "friday · understand + choose", "Concept lock within 90 minutes of the brief.", 8)
timeline = [
    ("5:00", "CAPTURE", "exact brief · rules · deliverables", PURPLE),
    ("5:30", "DECODE", "user · behaviour · constraint · success", CYAN),
    ("5:45", "DIVERGE", "six ideas each → cluster to three", LIME),
    ("6:10", "SCORE", "kill weakly testable + feature-heavy ideas", CORAL),
    ("6:30", "TALK", "five rapid user conversations", LILAC),
    ("7:25", "LOCK", "one line · one loop · one direction", PURPLE),
]
for i, (tm, act, detail, c) in enumerate(timeline):
    y = 1.55 + i * 0.82
    text(s, tm, 0.72, y + 0.08, 0.72, 0.25, 12, c, True)
    rect(s, 1.70, y + 0.07, 0.15, 0.15, c, radius=True)
    text(s, act, 2.15, y, 1.35, 0.32, 14, WHITE, True)
    text(s, detail, 3.68, y, 7.6, 0.34, 14, MUTED)
rect(s, 9.62, 5.15, 2.60, 1.28, PURPLE, radius=True)
text(s, "FRIDAY EXIT", 9.84, 5.39, 2.12, 0.22, 9, BG, True, align=PP_ALIGN.CENTER)
text(s, "everyone says the\nsame one-line product", 9.83, 5.72, 2.15, 0.50, 12, BG, True, align=PP_ALIGN.CENTER)

# 9 — Saturday
s = prs.slides.add_slide(blank)
add_bg(s); title(s, "saturday · prove + build", "Test low-fi early. Freeze features at 5pm.", 9)
blocks = [
    ("08:00", "MAP", "current journey + 5–7 frame hero loop", PURPLE),
    ("09:30", "LOW-FI", "clickable core", CYAN),
    ("10:30", "TEST", "3 users · observe, don’t explain", LIME),
    ("11:15", "MENTOR 1", "problem / solution logic", CORAL),
    ("12:00", "BUILD", "mid/high-fi + story in parallel", LILAC),
    ("14:00", "TEST 2", "comprehension + task success", PURPLE),
    ("15:30", "POLISH", "trust · error · accessibility state", CYAN),
    ("17:00", "FREEZE", "no new scope", LIME),
    ("18:15", "PITCH 1", "timed · brutal critique · backup demo", CORAL),
]
for i, (tm, head, body, c) in enumerate(blocks):
    col = 0 if i < 5 else 1
    row = i if i < 5 else i - 5
    x = 0.70 + col * 6.18
    y = 1.58 + row * 0.98
    rect(s, x, y, 5.76, 0.72, PANEL, radius=True, line=PANEL_2)
    text(s, tm, x + 0.20, y + 0.18, 0.78, 0.25, 11, c, True)
    text(s, head, x + 1.14, y + 0.17, 1.10, 0.25, 12, WHITE, True)
    text(s, body, x + 2.30, y + 0.16, 3.14, 0.32, 11.5, MUTED)
text(s, "Saturday exit: a stranger completes the loop, explains the value, and remembers the mechanism.", 0.72, 6.65, 11.9, 0.34, 13.5, WHITE, True)

# 10 — Sunday
s = prs.slides.add_slide(blank)
add_bg(s); title(s, "sunday · make it undeniable", "Submission safety first. Polish second.", 10)
card(s, 0.68, 1.65, 3.78, 2.0, "08:00–10:30", "STABILISE", "Fresh-eyes QA. Finish the deck and submission copy from the same story.", PURPLE)
card(s, 4.80, 1.65, 3.78, 2.0, "10:30–12:30", "ATTACK", "Ask a mentor: ‘Why would this not advance?’ Fix only decisive objections.", CORAL)
card(s, 8.92, 1.65, 3.78, 2.0, "12:30 onward", "REHEARSE", "Create a submittable checkpoint, then rehearse both 3- and 5-minute versions.", CYAN)
rect(s, 0.68, 4.17, 12.02, 1.38, PANEL, radius=True, line=PANEL_2)
pill(s, "-90 min", 0.96, 4.50, 1.10, LIME)
text(s, "A complete, tested submission exists before the official deadline.", 2.38, 4.48, 9.7, 0.35, 18, WHITE, True)
text(s, "Links work in incognito · assets disclosed · backup video on two devices · confirmation captured", 2.38, 4.97, 9.7, 0.34, 13, MUTED)
text(s, "After submission: freeze → hydrate → reset → present.", 0.70, 6.22, 11.8, 0.42, 21, LILAC, True)

# 11 — pitch
s = prs.slides.add_slide(blank)
add_bg(s); title(s, "the three-minute pitch", "Demo the changed behaviour — not a feature tour.", 11)
pitch = [
    ("0:00", "HUMAN MOMENT", "put the judge inside the pain", PURPLE),
    ("0:20", "EVIDENCE", "what people actually do", CYAN),
    ("0:45", "ONE LINE", "product + memorable mechanism", LIME),
    ("1:00", "HERO DEMO", "trigger → action → payoff", CORAL),
    ("1:50", "LEARNING", "what testing changed + trust", LILAC),
    ("2:15", "CREDIBILITY", "impact · feasibility · access", PURPLE),
    ("2:40", "CALLBACK", "return to the human; land one line", CYAN),
]
for i, (tm, head, body, c) in enumerate(pitch):
    x = 0.68 + i * 1.70
    rect(s, x, 1.73, 1.46, 3.74, PANEL, radius=True, line=PANEL_2)
    rect(s, x, 1.73, 1.46, 0.13, c, radius=True)
    text(s, tm, x + 0.13, 2.08, 1.18, 0.28, 13, c, True, align=PP_ALIGN.CENTER)
    text(s, head, x + 0.13, 2.62, 1.18, 0.65, 12, WHITE, True, align=PP_ALIGN.CENTER)
    text(s, body, x + 0.16, 3.48, 1.12, 1.10, 11, MUTED, align=PP_ALIGN.CENTER)
    if i < 6:
        text(s, "→", x + 1.45, 3.14, 0.24, 0.35, 16, MUTED, True, align=PP_ALIGN.CENTER)
pill(s, "speaker rule", 0.68, 6.10, 1.42, LILAC)
text(s, "One lead voice. Others own demo control and Q&A. Record a backup.", 2.34, 6.12, 9.8, 0.35, 15, WHITE, True)

# 12 — kill list
s = prs.slides.add_slide(blank)
add_bg(s); title(s, "the kill list", "Most teams lose through scope, not talent.", 12)
left = [
    "no ‘app for everyone’",
    "no more than 3 concepts after hour one",
    "no high-fi before low-fi testing",
    "no feature after Saturday 5pm",
    "no fake evidence or technical claims",
]
right = [
    "no vague ‘AI-powered’ magic",
    "no live demo without backup",
    "no five-person pixel committee",
    "no all-nighter as a strategy",
    "no process-heavy pitch",
]
bullet_list(s, left, 0.78, 1.75, 5.58, 3.95, 15, CORAL, 0.78)
bullet_list(s, right, 6.74, 1.75, 5.62, 3.95, 15, CORAL, 0.78)
rect(s, 0.70, 6.00, 11.95, 0.67, PURPLE, radius=True)
text(s, "Before adding anything: which score does this improve?", 0.90, 6.14, 11.55, 0.32, 18, BG, True, align=PP_ALIGN.CENTER)

# 13 — immediate actions
s = prs.slides.add_slide(blank)
add_bg(s); title(s, "the next 24 hours", "Win preparation starts with five boring certainties.", 13)
actions = [
    ("01", "REGISTRATION", "All five confirmation emails visible", PURPLE),
    ("02", "ROLE LOCK", "20-minute strengths audit tonight", CYAN),
    ("03", "ORGANISER CHECK", "brief · deadline · AI/IP · rubric/pitch", LIME),
    ("04", "90-MIN DRILL", "run the neutral sprint tomorrow", CORAL),
    ("05", "ARRIVAL", "Wilkinson G04 by 4:15pm Friday", LILAC),
]
for i, (num, head, body, c) in enumerate(actions):
    y = 1.55 + i * 0.97
    text(s, num, 0.78, y + 0.04, 0.50, 0.33, 13, c, True)
    rect(s, 1.55, y, 10.70, 0.70, PANEL, radius=True, line=PANEL_2)
    text(s, head, 1.86, y + 0.18, 2.25, 0.28, 13, WHITE, True)
    text(s, body, 4.33, y + 0.18, 7.55, 0.28, 13, MUTED)
text(s, "Target: the clearest tested idea in the room — delivered by the calmest team.", 0.74, 6.58, 11.75, 0.42, 19, LILAC, True, align=PP_ALIGN.CENTER)

# 14 — sources
s = prs.slides.add_slide(blank)
add_bg(s); title(s, "evidence + caveats", "Current organiser announcements override this deck.", 14)
sources = [
    ("2026 OFFICIAL", "events.humanitix.com/suede-designaton-2026", PURPLE),
    ("SUEDE LIVE LINKS", "linktr.ee/suedesociety", CYAN),
    ("2025 OVERVIEW", "suede-designathon-2025.devpost.com", LIME),
    ("2025 RULES", "suede-designathon-2025.devpost.com/rules", CORAL),
    ("2025 GALLERY", "suede-designathon-2025.devpost.com/project-gallery", LILAC),
    ("2024 RUBRIC", "suede-designathon-2024.devpost.com", PURPLE),
]
for i, (head, url, c) in enumerate(sources):
    y = 1.55 + i * 0.72
    text(s, head, 0.76, y + 0.04, 2.05, 0.25, 10, c, True)
    text(s, url, 3.05, y, 8.98, 0.32, 13.5, WHITE)
rect(s, 0.72, 6.15, 12.0, 0.64, PANEL, radius=True, line=PANEL_2)
text(s, "Research cut-off: 1 Sep 2026 · 2026 brief, rubric, deadline, judges, prizes, and AI/IP rules were not public.",
     0.94, 6.31, 11.55, 0.28, 12.5, MUTED, align=PP_ALIGN.CENTER)

prs.core_properties.title = "SUEDE Designathon 2026 — Winning Plan"
prs.core_properties.subject = "Competition operating proposal for Ethan, Darius, CC, Akari, and Chloe"
prs.core_properties.author = "OpenAI Codex for Ethan Wu"
prs.save(OUT)
print(OUT)
